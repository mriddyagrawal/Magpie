"""Shared file-content dispatch used by stage 1 (summarize.py) and stage 4 (answer.py).

Responsibilities:
- Know every supported file extension.
- Turn a single Path into a list of content blocks (strings + BinaryContent) suitable
  for handing to a PydanticAI Agent. Caller owns the framing (filename hints, question
  prefix, etc.); this module only produces the content itself.
"""

from __future__ import annotations

import re
from pathlib import Path

# `pydantic_ai` is deliberately NOT imported at module scope here. It pulls in
# the entire pydantic-ai agent graph (~800 ms cold-start cost as measured by
# `python -X importtime`), which would fire on every module that transitively
# imports `src.content` (src.answer, src.server, src.ingest.tier2, ...).
# The only thing we need from pydantic_ai is `BinaryContent`, used in two
# branches of `build_content_blocks` below — moved to a deferred import at
# function scope. See `Plans/Bundle Trim/Implementation Plan.md` PR-D.


IMAGE_EXTS = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
              ".webp": "image/webp", ".gif": "image/gif"}


# ---------------------------------------------------------------------------
# Index-time vision transcripts for scanned documents (2026-08-25 spike).
#
# The vision-isolation ladder showed the local VL model reads a SINGLE
# scanned page near-perfectly (GPAs, handwritten forms, a bilingual tax
# certificate) while the answer-time pixel path converted ~1/12 of scanned
# eval questions — same crowd/format story as text. So: transcribe pages
# ONCE at index time (the walker writes <APP_DATA_DIR>/transcripts/<key>.md
# through src/transcribe.py; Evaluations/transcribe_index.py is the same
# thing as a sweep over an already-indexed corpus), and at answer time read
# the transcript as TEXT — which also makes scanned docs answerable in cloud
# mode without pixels ever leaving the machine. No transcript → the pixel
# path below is unchanged.
# ---------------------------------------------------------------------------

def transcripts_dir() -> Path:
    """Where transcripts live. `MAGPIE_TRANSCRIPTS_DIR` overrides the default
    so two transcriber backends can be kept side by side and an eval arm can
    pick one without touching the other (Evaluations/transcript_recall.py)."""
    import os

    override = os.environ.get("MAGPIE_TRANSCRIPTS_DIR", "").strip()
    if override:
        return Path(override)
    from src.manifest import APP_DATA_DIR

    return APP_DATA_DIR / "transcripts"


def transcript_path_for(path: Path) -> Path:
    import hashlib

    key = hashlib.sha256(str(path.resolve()).lower().encode("utf-8")).hexdigest()[:16]
    return transcripts_dir() / f"{key}.md"


def transcript_for(path: Path) -> str | None:
    """The stored transcript for `path`, or None.

    A transcript with a header but no `## Page` body is a stub written for a
    picture that had no text in it (src/transcribe.py writes one so the file
    is not retried every sweep). It must read as "no transcript" here, or the
    reader would be handed a header instead of the pixels."""
    try:
        text = transcript_path_for(path).read_text(encoding="utf-8").strip()
    except OSError:
        return None
    if not text or "## Page" not in text:
        return None
    return text
CODE_EXTS = {".py", ".js", ".ts", ".tsx", ".jsx", ".go", ".rs", ".java",
             ".c", ".cpp", ".h", ".hpp", ".cs", ".rb", ".swift", ".kt",
             ".sh", ".sql", ".json", ".yaml", ".yml", ".toml"}
TEXT_EXTS = {".txt", ".log"}
MD_EXTS = {".md", ".markdown"}
CSV_EXTS = {".csv"}
PDF_EXTS = {".pdf"}
DOCX_EXTS = {".docx"}
XLSX_EXTS = {".xlsx", ".xlsm"}
PPTX_EXTS = {".pptx"}
HTML_EXTS = {".html", ".htm"}
IPYNB_EXTS = {".ipynb"}
ALT_EXTS = {".alt"}

PDF_VISION_DPI = 150

SUPPORTED_EXTS = (
    set(IMAGE_EXTS)
    | CSV_EXTS
    | PDF_EXTS
    | DOCX_EXTS
    | XLSX_EXTS
    | MD_EXTS
    | TEXT_EXTS
    | CODE_EXTS
    | PPTX_EXTS
    | HTML_EXTS
    | IPYNB_EXTS
)
# `.alt` is handled by `build_content_blocks` at answer time, but is deliberately
# NOT in SUPPORTED_EXTS — Stage 1's walker should skip .alt so it doesn't try to
# LLM-summarize them. Stage 3 owns .alt discovery and ingest.


class SummarizeError(RuntimeError):
    """Raised for per-file failures that should not abort a batch run."""


# Control characters that PDF extractors emit where a separator belonged.
# Stripping them is not cosmetic: `Receipt-2794-8324.pdf` in the sem_4 corpus
# extracts its invoice number as "9257BD07\x000001", and the consequences ran
# all the way down the pipeline — the summarizer produced a content-free
# summary AND leaked raw chat-template tokens (`<|tool_call_end|>`) into it,
# and four eval questions then came back parroting that empty summary. It was
# the only file in three corpora with NUL bytes, and the only one with leaked
# tokens. NUL is replaced with a space (it stands where a separator was);
# other C0 controls are dropped, keeping tab/newline/carriage return.
_CONTROL_CHARS = {
    ord(c): None
    for c in map(chr, list(range(0, 9)) + [11, 12] + list(range(14, 32)) + [127])
}
_CONTROL_CHARS[0] = " "


def scrub_control_chars(text: str) -> str:
    """Replace NUL with a space and drop other C0 control characters."""
    return text.translate(_CONTROL_CHARS)


def extract_pdf_text(path: Path, max_chars: int) -> str:
    from pypdf import PdfReader
    from pypdf.errors import PdfReadError

    try:
        reader = PdfReader(str(path))
    except (PdfReadError, OSError) as e:
        raise SummarizeError(f"could not open PDF {path}: {e}") from e

    if reader.is_encrypted:
        # Bank/government portal PDFs are routinely "encrypted" with an
        # EMPTY password — openable by anyone (pymupdf's needs_pass is
        # False for them). Treating is_encrypted as password-protected
        # made six corpus files invisible to indexing (visa-fee receipt,
        # Gettysburg ISFAAs — found 2026-08-26). Try the empty password;
        # only a PDF that actually refuses it is password-protected.
        try:
            decrypted = int(reader.decrypt(""))
        except Exception as e:  # noqa: BLE001 — any failure = truly locked
            raise SummarizeError(f"PDF is password-protected: {path}") from e
        if decrypted == 0:  # PasswordType.NOT_DECRYPTED
            raise SummarizeError(f"PDF is password-protected: {path}")
        # pypdf's AES content decryption needs the optional `cryptography`
        # backend we don't ship; pymupdf decrypts natively. Read
        # encrypted-but-open PDFs through pymupdf instead of adding a
        # dependency.
        return _extract_pdf_text_pymupdf(path, max_chars)

    chunks: list[str] = []
    total = 0

    # Bookmark-based TOC first. For long documents (books, textbooks, manuals)
    # the reader often asks "what's in this?" — the structured chapter list
    # answers that in a few hundred tokens, where sequential page extraction
    # would waste the whole budget on the preface.
    toc_block = _extract_pdf_toc(path)
    if toc_block:
        chunks.append(toc_block)
        total += len(toc_block)

    try:
        for page in reader.pages:
            if total >= max_chars:
                break
            t = page.extract_text() or ""
            chunks.append(t)
            total += len(t)
    except PdfReadError as e:
        raise SummarizeError(f"PDF parse error while reading pages: {path}: {e}") from e
    return "\n\n".join(chunks)


def _extract_pdf_text_pymupdf(path: Path, max_chars: int) -> str:
    """Text extraction via pymupdf, for PDFs pypdf cannot fully read
    (empty-password AES encryption). Same join format as the pypdf path;
    an empty result still means "scanned" to the caller."""
    try:
        import pymupdf
    except ImportError as e:
        raise SummarizeError(f"could not open PDF {path}: {e}") from e
    try:
        doc = pymupdf.open(str(path))
    except Exception as e:  # noqa: BLE001
        raise SummarizeError(f"could not open PDF {path}: {e}") from e
    try:
        if doc.needs_pass:
            raise SummarizeError(f"PDF is password-protected: {path}")
        chunks: list[str] = []
        total = 0
        for page in doc:
            if total >= max_chars:
                break
            t = page.get_text() or ""
            chunks.append(t)
            total += len(t)
    finally:
        doc.close()
    return "\n\n".join(chunks)


def _extract_pdf_toc(path: Path) -> str:
    """Return a formatted bookmark-based table of contents for a PDF, or "".

    Uses pymupdf's `doc.get_toc()` which reads the /Outlines object — the
    structured chapter index most publishers embed. Returns empty string if
    the PDF has no bookmarks (e.g. homemade scans, short receipts, etc.).
    """
    try:
        import pymupdf
    except ImportError:
        return ""
    try:
        doc = pymupdf.open(str(path))
    except Exception:  # pylint: disable=broad-except
        return ""
    try:
        toc = doc.get_toc() or []
    finally:
        doc.close()
    if not toc:
        return ""
    lines = ["## Table of Contents"]
    for level, title, _page in toc:
        indent = "  " * max(level - 1, 0)
        lines.append(f"{indent}- {title}")
    return "\n".join(lines)


_TOC_LINE_RE = re.compile(
    # Matches lines like:
    #   "1.2 Conservation of Momentum 47"
    #   "CHAPTER 3 Momentum and Angular Momentum 83"
    #   "Appendix A: Diagonalizing Real Symmetric Matrices 615"
    # The defining trait: ALL-CAPS or numeric prefix → title text → page number,
    # ending the line. A real content page rarely has 5+ lines of this shape.
    r"(?:^|\n)\s*(?:CHAPTER\s+|APPENDIX\s+|\d+(?:\.\d+)?)\s+"
    r"[A-Za-z][^\n]{2,80}?\s+\d+\s*(?=\n|$)",
    re.MULTILINE | re.IGNORECASE,
)

_BOOK_PAGE_LABEL_RE = re.compile(
    r"^\s*(\d{1,4}|[IVXLCDM]+|[ivxlcdm]+)\s*$",
    re.MULTILINE,
)


def _extract_book_page_label(text: str) -> str | None:
    """Recover the book's PRINTED page number from a PDF page's text.

    PDFs often carry the printed page number in the running header or footer.
    The PDF page index (1, 2, 3…) almost never matches the book's printed
    page (which restarts at 1 after front matter and may use roman numerals
    before that). Heuristic: look at the first 3 lines of the page (top-of-
    page page numbers — Taylor's textbook is one of these), then the last 3
    lines (footer style). A line is a page-label if it contains only a
    1-4 digit number or a Roman numeral, surrounded by whitespace.
    """
    if not text or not text.strip():
        return None
    lines = text.split("\n")
    for line in lines[:3]:
        m = _BOOK_PAGE_LABEL_RE.match(line)
        if m:
            return m.group(1)
    for line in reversed(lines[-3:]):
        m = _BOOK_PAGE_LABEL_RE.match(line)
        if m:
            return m.group(1)
    return None


def _looks_like_toc_page(text: str) -> bool:
    """True when the page is dominated by table-of-contents-shaped lines.

    TOC pages list every chapter and section by name → they outscore real
    content pages on any "what's in this textbook" query. The bookmark TOC
    is already in the summary tier; we don't want it doubled up here.
    """
    matches = _TOC_LINE_RE.findall(text)
    return len(matches) >= 5


def extract_pdf_relevant_pages(
    path: Path,
    keywords: list[str],
    *,
    max_pages: int = 10,
    max_chars: int = 25_000,
    context_pages: int = 1,
) -> str:
    """Pick the PDF pages most relevant to `keywords`; return their text.

    The lazy-chunking primitive used at answer time for long PDFs that would
    otherwise have their first ~25K chars (cover + preface) cut and sent to
    the LLM. Algorithm: score each page by case-insensitive keyword hits,
    take top `max_pages`, expand by `context_pages` neighbors, emit in
    document order with `## PDF page N (book p. X)` anchors capped at
    `max_chars`. Returns "" when keywords empty, no page matches, or PDF
    can't be opened — caller falls back to first-N-chars extract.
    """
    if not keywords:
        return ""
    try:
        import pymupdf
    except ImportError:
        return ""
    try:
        doc = pymupdf.open(str(path))
    except Exception:  # pylint: disable=broad-except
        return ""

    try:
        n_pages = len(doc)
        if n_pages == 0:
            return ""
        kws = [k.lower() for k in keywords if k.strip()]
        if not kws:
            return ""

        page_scores: list[tuple[int, int]] = []
        for i in range(n_pages):
            try:
                raw_text = doc[i].get_text() or ""
            except Exception:  # pylint: disable=broad-except
                continue
            if _looks_like_toc_page(raw_text):
                continue
            text = raw_text.lower()
            score = sum(text.count(kw) for kw in kws)
            if score > 0:
                page_scores.append((i, score))

        if not page_scores:
            return ""

        page_scores.sort(key=lambda kv: -kv[1])
        seeds = {p for p, _ in page_scores[:max_pages]}
        with_context: set[int] = set()
        for p in seeds:
            for delta in range(-context_pages, context_pages + 1):
                if 0 <= p + delta < n_pages:
                    with_context.add(p + delta)

        chunks: list[str] = []
        total = 0
        for i in sorted(with_context):
            try:
                page_text = doc[i].get_text() or ""
            except Exception:  # pylint: disable=broad-except
                continue
            book_label = _extract_book_page_label(page_text)
            if book_label:
                anchor = f"## PDF page {i + 1} (book p. {book_label})"
            else:
                anchor = f"## PDF page {i + 1}"
            block = f"{anchor}\n{page_text}"
            chunks.append(block)
            total += len(block)
            if total >= max_chars:
                break

        return "\n\n".join(chunks)
    finally:
        doc.close()


def render_pdf_pages_as_png(path: Path, max_pages: int) -> list[bytes]:
    import pymupdf

    try:
        doc = pymupdf.open(str(path))
    except Exception as e:
        raise SummarizeError(f"could not open PDF for rendering: {path}: {e}") from e

    images: list[bytes] = []
    try:
        if doc.needs_pass:
            raise SummarizeError(f"PDF is password-protected: {path}")
        for i, page in enumerate(doc):
            if i >= max_pages:
                break
            pix = page.get_pixmap(dpi=PDF_VISION_DPI)
            images.append(pix.tobytes("png"))
    finally:
        doc.close()
    return images


def extract_docx_text(path: Path) -> str:
    from docx import Document
    from docx.opc.exceptions import PackageNotFoundError

    try:
        doc = Document(str(path))
    except (PackageNotFoundError, OSError) as e:
        raise SummarizeError(f"not a valid .docx file: {path}: {e}") from e

    parts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append("\t".join(cells))
    return "\n".join(parts)


def extract_pptx_text(path: Path) -> str:
    """Extract slide titles + bullet text + speaker notes from a .pptx.

    Returns concatenated text with one "## Slide N" header per slide so
    retrieval hits carry slide context. Skips images / charts / SmartArt —
    text only. Image-heavy decks are expected to be additionally routed
    through T4 (ColPali) for visual retrieval.
    """
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError

    try:
        prs = Presentation(str(path))
    except (PackageNotFoundError, OSError) as e:
        raise SummarizeError(f"not a valid .pptx file: {path}: {e}") from e

    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_bits: list[str] = [f"## Slide {i}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_bits.append(shape.text)
        notes = slide.notes_slide if slide.has_notes_slide else None
        if notes is not None:
            note_text = notes.notes_text_frame.text.strip() if notes.notes_text_frame else ""
            if note_text:
                slide_bits.append(f"[notes] {note_text}")
        if len(slide_bits) > 1:
            parts.append("\n".join(slide_bits))
    return "\n\n".join(parts)


def extract_html_text(path: Path) -> str:
    """Extract clean article text from an .html/.htm file via trafilatura.

    Strips boilerplate (nav, header/footer, ad chrome, script/style). If
    extraction yields no usable content (e.g. JS-only SPAs that render
    nothing on the server), raises SummarizeError.
    """
    import trafilatura

    try:
        raw = path.read_bytes()
    except OSError as e:
        raise SummarizeError(f"could not read html {path}: {e}") from e

    text = trafilatura.extract(raw, include_comments=False, include_tables=True) or ""
    text = text.strip()
    if not text:
        # Fallback: lossy UTF-8 decode of the raw bytes, so SPAs don't
        # silently vanish — the user at least gets whatever static content
        # the HTML carried, tags and all.
        try:
            text = raw.decode("utf-8", errors="ignore").strip()
        except Exception:  # pylint: disable=broad-except
            text = ""
    if not text:
        raise SummarizeError(f"html extracted empty text: {path}")
    return text


def extract_ipynb_text(path: Path) -> str:
    """Extract code + markdown cell sources from a Jupyter notebook.

    Skips cell outputs (often noisy, sometimes binary). Cell boundaries are
    surfaced with a "# Cell N (type)" header so retrieval hits know which
    cell a match came from.
    """
    import json as _json

    try:
        raw = path.read_text(encoding="utf-8")
    except (UnicodeDecodeError, OSError) as e:
        raise SummarizeError(f"could not read ipynb {path}: {e}") from e
    try:
        nb = _json.loads(raw)
    except _json.JSONDecodeError as e:
        raise SummarizeError(f"ipynb is not valid JSON: {path}: {e}") from e

    cells = nb.get("cells", [])
    if not isinstance(cells, list):
        raise SummarizeError(f"ipynb has no cells list: {path}")

    parts: list[str] = []
    for i, cell in enumerate(cells, start=1):
        ctype = cell.get("cell_type", "unknown")
        source = cell.get("source", "")
        if isinstance(source, list):
            source = "".join(source)
        source = source.strip()
        if source:
            parts.append(f"# Cell {i} ({ctype})\n{source}")
    return "\n\n".join(parts)


def extract_xlsx_text(path: Path) -> str:
    from openpyxl import load_workbook
    from openpyxl.utils.exceptions import InvalidFileException

    try:
        wb = load_workbook(str(path), data_only=True, read_only=True)
    except (InvalidFileException, OSError) as e:
        raise SummarizeError(f"not a valid .xlsx file: {path}: {e}") from e

    parts: list[str] = []
    try:
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"## Sheet: {sheet_name}")
            for row in ws.iter_rows(values_only=True):
                if not any(cell is not None and str(cell).strip() for cell in row):
                    continue
                parts.append(",".join("" if cell is None else str(cell) for cell in row))
            parts.append("")
    finally:
        wb.close()
    return "\n".join(parts)


def build_content_blocks(
    path: Path,
    *,
    max_chars: int,
    max_pdf_pages: int,
    search_keywords: list[str] | None = None,
) -> list:
    """Scrubbing wrapper — see `_build_content_blocks`.

    Every textual block leaves through here, whatever produced it, so control
    characters are stripped in exactly one place instead of at a dozen
    `return` statements.
    """
    return [
        scrub_control_chars(b) if isinstance(b, str) else b
        for b in _build_content_blocks(
            path,
            max_chars=max_chars,
            max_pdf_pages=max_pdf_pages,
            search_keywords=search_keywords,
        )
    ]


def _build_content_blocks(
    path: Path,
    *,
    max_chars: int,
    max_pdf_pages: int,
    search_keywords: list[str] | None = None,
) -> list:
    """Return content blocks for a single file.

    The returned list contains `str` blocks (carrying textual content, with a
    'Content type: ...' header) and `BinaryContent` blocks (for images or for PDF
    pages rendered via the scanned-PDF vision fallback).

    `search_keywords` (optional): if provided AND the PDF has more text than
    `max_chars` would fit, the lazy-chunking primitive
    `extract_pdf_relevant_pages` picks the matching pages first and falls
    back to the front-of-file extract if no page matches. For short PDFs
    that fit entirely in `max_chars`, the keywords are ignored — the whole
    file goes to the LLM regardless.

    Raises SummarizeError for unsupported extensions, encrypted PDFs, corrupt
    Office files, non-UTF-8 text files, or empty content.
    """
    # Deferred — only paid when this function is actually called, not on
    # every import of src.content. See module-level comment.
    from pydantic_ai import BinaryContent

    ext = path.suffix.lower()

    if ext in IMAGE_EXTS:
        # A photographed receipt is the same object as a one-page scanned PDF:
        # pixels of a document. Read its index-time transcript as text when
        # one exists (same reasons as the PDF branch below — cacheable,
        # groundable, cloud-safe); fall back to the pixels only when it does
        # not. Until 2026-08-29 images skipped this check and always went to
        # the reader as pixels, so the transcript work never reached them.
        transcript = transcript_for(path)
        if transcript:
            return [
                "Content type: image (reading the index-time transcript)"
                f"\n\n---\n{transcript[:max_chars]}"
            ]
        return [BinaryContent(data=path.read_bytes(), media_type=IMAGE_EXTS[ext])]

    if ext in PDF_EXTS:
        # Always extract first to know the file's full size.
        full_text = extract_pdf_text(path, max_chars * 4).strip()

        if full_text:
            # Lazy-chunking trigger: file is bigger than the LLM budget AND
            # caller gave us keywords. Otherwise the cheap front-of-file
            # extract is fine — short PDFs fit entirely already.
            if search_keywords and len(full_text) > max_chars:
                relevant = extract_pdf_relevant_pages(
                    path, keywords=search_keywords, max_chars=max_chars
                )
                if relevant:
                    return [
                        "Content type: pdf (long file — pages selected by query "
                        "keywords; see ## Page N anchors below)\n\n---\n"
                        f"{relevant[:max_chars]}"
                    ]
                # No keyword matched — fall through to front-of-file.
            return [f"Content type: pdf\n\n---\n{full_text[:max_chars]}"]

        # Empty extract → scanned/image-only PDF. Prefer an index-time
        # vision transcript (see transcript_for above): text path, cheaper,
        # measured far more accurate than answer-time pixel reading.
        transcript = transcript_for(path)
        if transcript:
            return [
                "Content type: pdf (scanned — reading the index-time vision "
                f"transcript)\n\n---\n{transcript[:max_chars]}"
            ]
        # No transcript yet: render pages as images (unchanged fallback).
        pages = render_pdf_pages_as_png(path, max_pdf_pages)
        if not pages:
            raise SummarizeError(f"PDF has no pages: {path}")
        blocks: list = [
            f"Content type: pdf (scanned / image-only — {len(pages)} page(s) as images)"
        ]
        for page_png in pages:
            blocks.append(BinaryContent(data=page_png, media_type="image/png"))
        return blocks

    if ext in DOCX_EXTS:
        text = extract_docx_text(path).strip()
        if not text:
            raise SummarizeError(f"docx appears empty: {path}")
        return [f"Content type: docx\n\n---\n{text[:max_chars]}"]

    if ext in XLSX_EXTS:
        text = extract_xlsx_text(path).strip()
        if not text:
            raise SummarizeError(f"xlsx appears empty: {path}")
        return [f"Content type: xlsx\n\n---\n{text[:max_chars]}"]

    if ext in PPTX_EXTS:
        text = extract_pptx_text(path).strip()
        if not text:
            raise SummarizeError(f"pptx appears empty: {path}")
        return [f"Content type: pptx\n\n---\n{text[:max_chars]}"]

    if ext in HTML_EXTS:
        text = extract_html_text(path).strip()
        return [f"Content type: html\n\n---\n{text[:max_chars]}"]

    if ext in IPYNB_EXTS:
        text = extract_ipynb_text(path).strip()
        if not text:
            raise SummarizeError(f"ipynb appears empty: {path}")
        return [f"Content type: ipynb\n\n---\n{text[:max_chars]}"]

    if ext in CSV_EXTS:
        try:
            text = path.read_text(encoding="utf-8")
        except UnicodeDecodeError as e:
            raise SummarizeError(f"file is not valid UTF-8 text: {path}") from e
        if not text.strip():
            raise SummarizeError(f"csv appears empty: {path}")
        return [f"Content type: csv\n\n---\n{text[:max_chars]}"]

    if ext in MD_EXTS or ext in TEXT_EXTS or ext in CODE_EXTS:
        ctype = "markdown" if ext in MD_EXTS else ("code" if ext in CODE_EXTS else "text")
        try:
            text = path.read_text(encoding="utf-8")
        except UnicodeDecodeError as e:
            raise SummarizeError(f"file is not valid UTF-8 text: {path}") from e
        return [f"Content type: {ctype}\n\n---\n{text[:max_chars]}"]

    if ext in ALT_EXTS:
        try:
            text = path.read_text(encoding="utf-8")
        except UnicodeDecodeError as e:
            raise SummarizeError(f"file is not valid UTF-8 text: {path}") from e
        if not text.strip():
            raise SummarizeError(f"alt file appears empty: {path}")
        return [
            "Content type: alt (structured sidecar for a too-large source "
            "file — treat this YAML as the authoritative content; the "
            "referenced binary is not loaded).\n\n---\n" + text[:max_chars]
        ]

    raise SummarizeError(f"unsupported file type '{ext}' for {path}")
