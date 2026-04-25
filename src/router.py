"""Ingest router: peek each file, score it, decide which tier to run.

See `Plans/Indexing Tiers.md` for the full policy. This module is pure —
no side effects, no database, no network. Callers are responsible for
actually running the chosen tier.

The router's contract:

    peek_result = peek(path)
    decision = decide(
        peek_result,
        nasconfig=load_nasconfig(path),
        gpu_available=...,
        t4_budget_used_mb=...,
    )
    # decision.routes is a list like ["T3", "T2"] or ["T4"] or ["T1"]
    # decision.skip_reason is set if the file was rejected

Extensions to "supported" are driven entirely by the tier-decision logic;
there is no global SUPPORTED_EXTS. Unknown extensions route to "skip" with
a clear reason.
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Literal

Tier = Literal["T0", "T1", "T2", "T3", "T4"]
Criticality = Literal["critical", "normal", "casual"]
CriticalitySource = Literal["user", "auto", "default"]


# ---------------------------------------------------------------------------
# File-type buckets (extension only — routing goes deeper with peek)
# ---------------------------------------------------------------------------

TEXT_EXTS = {".txt", ".md", ".markdown", ".log"}
CODE_EXTS = {
    ".py", ".js", ".ts", ".tsx", ".jsx", ".go", ".rs", ".java",
    ".c", ".cpp", ".h", ".hpp", ".cs", ".rb", ".swift", ".kt",
    ".sh", ".sql",
}
CONFIG_EXTS = {".json", ".yaml", ".yml", ".toml"}
CSV_EXTS = {".csv"}
PDF_EXTS = {".pdf"}
DOCX_EXTS = {".docx"}
XLSX_EXTS = {".xlsx", ".xlsm"}
PPTX_EXTS = {".pptx"}
HTML_EXTS = {".html", ".htm"}
IPYNB_EXTS = {".ipynb"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".gif"}


# ---------------------------------------------------------------------------
# Size and count thresholds (documented in Plans/Indexing Tiers.md)
# ---------------------------------------------------------------------------

TEXT_SIZE_T0_THRESHOLD = 100 * 1024          # 100 KB
CODE_SIZE_T0_THRESHOLD = 500 * 1024          # 500 KB
CONFIG_SIZE_T0_THRESHOLD = 50 * 1024         # 50 KB

CSV_ROWS_T1_MAX = 1_000
CSV_ROWS_T2_MAX = 100_000

PDF_SHORT_PAGE_THRESHOLD = 5                 # ≤5 pages → typically discriminator-heavy

# Image thumbnail skip thresholds — small files are UI assets, not docs. The
# size+dim rule uses AND (both must be below the threshold to skip), so real
# document scans (large bytes even if low-res, or full-res but aggressively
# compressed) pass through. Raised from 200→600 px on 2026-04-21 after finding
# that ColPali was encoding 500+ full-res stock photos as "documents" because
# they passed the old 200 px floor. Normal document scans are ≥1200 px in at
# least one dim; decorative clip-art is typically ≤500 px.
IMAGE_THUMBNAIL_SIZE_BYTES = 50 * 1024
IMAGE_THUMBNAIL_MIN_DIM = 600                # px

# T4 cost estimates (per page)
T4_STORAGE_MB_PER_PAGE = 0.2                 # int8-quantized multi-vector
T4_GPU_SECONDS_PER_PAGE = 1.0
T4_CPU_SECONDS_PER_PAGE = 10.0

# T4 per-file gates
T4_MAX_STORAGE_MB_PER_FILE = 50.0
T4_MAX_SECONDS_PER_FILE_GPU = 30.0
T4_MAX_SECONDS_PER_FILE_CPU = 10.0

# Default corpus-wide T4 storage budget (MB). Overridable via .nasconfig.yaml.
DEFAULT_T4_BUDGET_MB = 5 * 1024              # 5 GB

# Size of peek-text slice used for sensitivity scoring.
PEEK_TEXT_MAX_CHARS = 5000


# ---------------------------------------------------------------------------
# PeekResult — the cheap-to-compute inspection output
# ---------------------------------------------------------------------------

@dataclass(frozen=True)
class PeekResult:
    """Deterministic inspection result for one file. All fields populated by `peek()`."""

    path: Path
    ext: str                          # lowercased suffix including leading dot
    size_bytes: int
    # PDF / DOCX text shape:
    page_count: int                   # 0 for non-paginated types
    text_density: int                 # chars per sampled page (PDF/DOCX) or chars per file (text)
    extractable: bool                 # did extraction yield real-word content
    # DOCX layout signal:
    image_ratio: float                # #images / #paragraphs, 0.0 for non-DOCX
    # CSV:
    row_count: int | None             # None for non-CSV
    # Image:
    image_dims: tuple[int, int] | None  # (w, h) in px, None if not an image
    # Small text slice used for sensitivity scoring (UTF-8, up to PEEK_TEXT_MAX_CHARS)
    peek_text: str
    peek_error: str | None = None


# ---------------------------------------------------------------------------
# RouteDecision — what the router tells the caller to do
# ---------------------------------------------------------------------------

@dataclass(frozen=True)
class RouteDecision:
    """One file's dispatch verdict. `routes` is the list of tier workers to run."""

    routes: list[Tier]
    visual_score: int
    sensitivity_score: int
    t4_cost_mb: float
    t4_cost_s: float
    criticality: Criticality
    criticality_source: CriticalitySource
    skip_reason: str | None = None            # non-None means: don't ingest
    notes: list[str] = field(default_factory=list)

    @property
    def skipped(self) -> bool:
        return self.skip_reason is not None


# ---------------------------------------------------------------------------
# Peek helpers (deferred imports — this module must stay import-cheap)
# ---------------------------------------------------------------------------

def _read_bytes_safely(path: Path, max_bytes: int) -> tuple[bytes, str | None]:
    try:
        with path.open("rb") as f:
            return f.read(max_bytes), None
    except OSError as e:
        return b"", f"read failed: {e}"


_EXTRACTABLE_MIN_LETTERS = 5


def _decode_peek(data: bytes) -> tuple[str, bool]:
    """Decode bytes to UTF-8 text; return (text, extractable). Garbage → ("", False).

    "extractable" means: we got at least ~20 alphabetic characters AND letters
    make up at least 30% of non-whitespace content. An empty / whitespace-only
    peek is NOT extractable — that's how we detect scanned PDFs whose text
    extraction yielded nothing usable.
    """
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError:
        try:
            text = data.decode("utf-8", errors="ignore")
        except Exception:  # pylint: disable=broad-except
            return "", False

    if not text:
        return "", False

    letters = sum(1 for c in text if c.isalpha())
    if letters < _EXTRACTABLE_MIN_LETTERS:
        return text[:PEEK_TEXT_MAX_CHARS], False

    non_ws = sum(1 for c in text if not c.isspace())
    if non_ws == 0:
        return text[:PEEK_TEXT_MAX_CHARS], False
    ratio = letters / non_ws
    return text[:PEEK_TEXT_MAX_CHARS], ratio >= 0.3


def _peek_text_file(path: Path) -> PeekResult:
    size = path.stat().st_size
    data, err = _read_bytes_safely(path, PEEK_TEXT_MAX_CHARS * 2)
    text, extractable = _decode_peek(data)
    return PeekResult(
        path=path,
        ext=path.suffix.lower(),
        size_bytes=size,
        page_count=0,
        text_density=len(text),
        extractable=extractable,
        image_ratio=0.0,
        row_count=None,
        image_dims=None,
        peek_text=text,
        peek_error=err,
    )


def _peek_csv(path: Path) -> PeekResult:
    """Count CSV rows (excluding header) streamingly; take first chunk as peek_text."""
    size = path.stat().st_size
    row_count = 0
    peek_text = ""
    err: str | None = None
    try:
        with path.open(encoding="utf-8", errors="ignore") as f:
            header = f.readline()
            peek_lines: list[str] = [header] if header else []
            for line in f:
                row_count += 1
                if sum(len(x) for x in peek_lines) < PEEK_TEXT_MAX_CHARS:
                    peek_lines.append(line)
            peek_text = "".join(peek_lines)[:PEEK_TEXT_MAX_CHARS]
    except OSError as e:
        err = f"read failed: {e}"

    return PeekResult(
        path=path,
        ext=path.suffix.lower(),
        size_bytes=size,
        page_count=0,
        text_density=len(peek_text),
        extractable=bool(peek_text),
        image_ratio=0.0,
        row_count=row_count,
        image_dims=None,
        peek_text=peek_text,
        peek_error=err,
    )


def _peek_pdf(path: Path) -> PeekResult:
    """Three-point sample (first/middle/last) of embedded text; return density + count."""
    size = path.stat().st_size
    err: str | None = None
    page_count = 0
    total_chars = 0
    peek_text = ""
    try:
        import pymupdf  # deferred: heavy import
    except ImportError as e:
        return PeekResult(
            path=path, ext=path.suffix.lower(), size_bytes=size,
            page_count=0, text_density=0, extractable=False,
            image_ratio=0.0, row_count=None, image_dims=None,
            peek_text="", peek_error=f"pymupdf unavailable: {e}",
        )

    try:
        doc = pymupdf.open(str(path))
    except Exception as e:  # pylint: disable=broad-except
        return PeekResult(
            path=path, ext=path.suffix.lower(), size_bytes=size,
            page_count=0, text_density=0, extractable=False,
            image_ratio=0.0, row_count=None, image_dims=None,
            peek_text="", peek_error=f"pdf open failed: {e}",
        )

    try:
        page_count = len(doc)
        if page_count == 0:
            err = "pdf has 0 pages"
        else:
            # Three-point sample.
            if page_count <= 3:
                idxs = list(range(page_count))
            else:
                idxs = sorted({0, page_count // 2, page_count - 1})

            chunks: list[str] = []
            for i in idxs:
                try:
                    t = doc[i].get_text() or ""
                except Exception:  # pylint: disable=broad-except
                    t = ""
                total_chars += len(t)
                chunks.append(t)
            peek_text = "\n".join(chunks)[:PEEK_TEXT_MAX_CHARS]
            avg_chars = total_chars // max(len(idxs), 1)
            total_chars = avg_chars  # overload: text_density = chars/page
    finally:
        doc.close()

    _, extractable = _decode_peek(peek_text.encode("utf-8"))
    return PeekResult(
        path=path,
        ext=path.suffix.lower(),
        size_bytes=size,
        page_count=page_count,
        text_density=total_chars,
        extractable=extractable,
        image_ratio=0.0,
        row_count=None,
        image_dims=None,
        peek_text=peek_text,
        peek_error=err,
    )


def _peek_docx(path: Path) -> PeekResult:
    size = path.stat().st_size
    err: str | None = None
    n_para = 0
    n_image = 0
    peek_text = ""
    try:
        from docx import Document  # deferred
        from docx.opc.exceptions import PackageNotFoundError
    except ImportError as e:
        return PeekResult(
            path=path, ext=path.suffix.lower(), size_bytes=size,
            page_count=0, text_density=0, extractable=False,
            image_ratio=0.0, row_count=None, image_dims=None,
            peek_text="", peek_error=f"python-docx unavailable: {e}",
        )

    try:
        doc = Document(str(path))
    except (PackageNotFoundError, OSError) as e:
        return PeekResult(
            path=path, ext=path.suffix.lower(), size_bytes=size,
            page_count=0, text_density=0, extractable=False,
            image_ratio=0.0, row_count=None, image_dims=None,
            peek_text="", peek_error=f"docx open failed: {e}",
        )

    pieces: list[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            n_para += 1
            if sum(len(x) for x in pieces) < PEEK_TEXT_MAX_CHARS:
                pieces.append(p.text)
    # Embedded image parts live in the package's relationships.
    try:
        rels = doc.part.rels
        n_image = sum(1 for r in rels.values() if "image" in str(r.reltype).lower())
    except Exception:  # pylint: disable=broad-except
        n_image = 0

    peek_text = "\n".join(pieces)[:PEEK_TEXT_MAX_CHARS]
    image_ratio = n_image / n_para if n_para > 0 else (1.0 if n_image > 0 else 0.0)
    _, extractable = _decode_peek(peek_text.encode("utf-8"))

    return PeekResult(
        path=path,
        ext=path.suffix.lower(),
        size_bytes=size,
        page_count=max(n_para // 30, 1),     # rough: ~30 paragraphs per page
        text_density=len(peek_text),          # total peek chars, not per-page for DOCX
        extractable=extractable,
        image_ratio=image_ratio,
        row_count=None,
        image_dims=None,
        peek_text=peek_text,
        peek_error=err,
    )


def _peek_xlsx(path: Path) -> PeekResult:
    """Coarse peek: count sheets, take first N cells as peek_text."""
    size = path.stat().st_size
    err: str | None = None
    peek_text = ""
    n_sheets = 0
    try:
        from openpyxl import load_workbook  # deferred
        from openpyxl.utils.exceptions import InvalidFileException
    except ImportError as e:
        return PeekResult(
            path=path, ext=path.suffix.lower(), size_bytes=size,
            page_count=0, text_density=0, extractable=False,
            image_ratio=0.0, row_count=None, image_dims=None,
            peek_text="", peek_error=f"openpyxl unavailable: {e}",
        )

    try:
        wb = load_workbook(str(path), data_only=True, read_only=True)
    except (InvalidFileException, OSError) as e:
        return PeekResult(
            path=path, ext=path.suffix.lower(), size_bytes=size,
            page_count=0, text_density=0, extractable=False,
            image_ratio=0.0, row_count=None, image_dims=None,
            peek_text="", peek_error=f"xlsx open failed: {e}",
        )

    pieces: list[str] = []
    try:
        for name in wb.sheetnames:
            n_sheets += 1
            ws = wb[name]
            for row in ws.iter_rows(values_only=True):
                line = ",".join("" if c is None else str(c) for c in row)
                if line.strip():
                    pieces.append(line)
                if sum(len(p) for p in pieces) >= PEEK_TEXT_MAX_CHARS:
                    break
            if sum(len(p) for p in pieces) >= PEEK_TEXT_MAX_CHARS:
                break
    finally:
        wb.close()

    peek_text = "\n".join(pieces)[:PEEK_TEXT_MAX_CHARS]
    _, extractable = _decode_peek(peek_text.encode("utf-8"))
    return PeekResult(
        path=path,
        ext=path.suffix.lower(),
        size_bytes=size,
        page_count=n_sheets,
        text_density=len(peek_text),
        extractable=extractable,
        image_ratio=0.0,
        row_count=None,
        image_dims=None,
        peek_text=peek_text,
        peek_error=err,
    )


def _peek_pptx(path: Path) -> PeekResult:
    """Count slides and image-to-paragraph-ish ratio. Pulls a text sample."""
    size = path.stat().st_size
    err: str | None = None
    n_slides = 0
    n_image_shapes = 0
    n_text_paragraphs = 0
    pieces: list[str] = []
    try:
        from pptx import Presentation  # deferred
        from pptx.exc import PackageNotFoundError
    except ImportError as e:
        return PeekResult(
            path=path, ext=path.suffix.lower(), size_bytes=size,
            page_count=0, text_density=0, extractable=False,
            image_ratio=0.0, row_count=None, image_dims=None,
            peek_text="", peek_error=f"python-pptx unavailable: {e}",
        )
    try:
        prs = Presentation(str(path))
    except (PackageNotFoundError, OSError) as e:
        return PeekResult(
            path=path, ext=path.suffix.lower(), size_bytes=size,
            page_count=0, text_density=0, extractable=False,
            image_ratio=0.0, row_count=None, image_dims=None,
            peek_text="", peek_error=f"pptx open failed: {e}",
        )

    for slide in prs.slides:
        n_slides += 1
        for shape in slide.shapes:
            # Picture shapes have shape_type == 13 (MSO_SHAPE_TYPE.PICTURE)
            if getattr(shape, "shape_type", None) == 13:
                n_image_shapes += 1
            text = getattr(shape, "text", "") or ""
            if text.strip():
                n_text_paragraphs += 1
                if sum(len(p) for p in pieces) < PEEK_TEXT_MAX_CHARS:
                    pieces.append(text)
    peek_text = "\n".join(pieces)[:PEEK_TEXT_MAX_CHARS]
    image_ratio = (
        n_image_shapes / n_text_paragraphs if n_text_paragraphs > 0
        else (1.0 if n_image_shapes > 0 else 0.0)
    )
    _, extractable = _decode_peek(peek_text.encode("utf-8"))
    return PeekResult(
        path=path, ext=path.suffix.lower(), size_bytes=size,
        page_count=n_slides,
        text_density=len(peek_text) // max(n_slides, 1),
        extractable=extractable,
        image_ratio=image_ratio,
        row_count=None, image_dims=None,
        peek_text=peek_text, peek_error=err,
    )


def _peek_html(path: Path) -> PeekResult:
    """Extract clean text via trafilatura; report its length + extractability."""
    size = path.stat().st_size
    err: str | None = None
    clean = ""
    try:
        import trafilatura  # deferred
    except ImportError as e:
        return PeekResult(
            path=path, ext=path.suffix.lower(), size_bytes=size,
            page_count=0, text_density=0, extractable=False,
            image_ratio=0.0, row_count=None, image_dims=None,
            peek_text="", peek_error=f"trafilatura unavailable: {e}",
        )
    try:
        raw = path.read_bytes()
        clean = trafilatura.extract(raw, include_comments=False) or ""
    except OSError as e:
        err = f"html read failed: {e}"
    except Exception as e:  # pylint: disable=broad-except
        err = f"trafilatura failed: {type(e).__name__}: {e}"

    clean = clean.strip()
    peek_text = clean[:PEEK_TEXT_MAX_CHARS]
    _, extractable = _decode_peek(peek_text.encode("utf-8"))
    return PeekResult(
        path=path, ext=path.suffix.lower(), size_bytes=size,
        page_count=1, text_density=len(peek_text),
        extractable=extractable, image_ratio=0.0,
        row_count=None, image_dims=None,
        peek_text=peek_text, peek_error=err,
    )


def _peek_ipynb(path: Path) -> PeekResult:
    """Count cells, extract a text sample from the first few."""
    import json as _json
    size = path.stat().st_size
    err: str | None = None
    cells: list = []
    pieces: list[str] = []
    try:
        raw = path.read_text(encoding="utf-8")
        nb = _json.loads(raw)
        cells = nb.get("cells", []) if isinstance(nb, dict) else []
    except (OSError, UnicodeDecodeError) as e:
        err = f"ipynb read failed: {e}"
    except _json.JSONDecodeError as e:
        err = f"ipynb is not valid JSON: {e}"

    for cell in cells:
        if not isinstance(cell, dict):
            continue
        source = cell.get("source", "")
        if isinstance(source, list):
            source = "".join(source)
        if isinstance(source, str) and source.strip():
            pieces.append(source)
        if sum(len(p) for p in pieces) >= PEEK_TEXT_MAX_CHARS:
            break
    peek_text = "\n".join(pieces)[:PEEK_TEXT_MAX_CHARS]
    _, extractable = _decode_peek(peek_text.encode("utf-8"))
    return PeekResult(
        path=path, ext=path.suffix.lower(), size_bytes=size,
        page_count=len(cells),
        text_density=len(peek_text) // max(len(cells), 1),
        extractable=extractable, image_ratio=0.0,
        row_count=None, image_dims=None,
        peek_text=peek_text, peek_error=err,
    )


def _peek_image(path: Path) -> PeekResult:
    size = path.stat().st_size
    dims: tuple[int, int] | None = None
    err: str | None = None
    try:
        from PIL import Image  # deferred
    except ImportError as e:
        return PeekResult(
            path=path, ext=path.suffix.lower(), size_bytes=size,
            page_count=1, text_density=0, extractable=False,
            image_ratio=1.0, row_count=None, image_dims=None,
            peek_text="", peek_error=f"PIL unavailable: {e}",
        )
    try:
        with Image.open(str(path)) as im:
            dims = im.size
    except Exception as e:  # pylint: disable=broad-except
        err = f"image open failed: {e}"

    return PeekResult(
        path=path,
        ext=path.suffix.lower(),
        size_bytes=size,
        page_count=1,
        text_density=0,
        extractable=False,
        image_ratio=1.0,
        row_count=None,
        image_dims=dims,
        peek_text="",
        peek_error=err,
    )


def peek(path: Path) -> PeekResult:
    """Dispatch to the right peeker based on extension. Pure-ish: does read I/O only."""
    ext = path.suffix.lower()

    if ext in TEXT_EXTS or ext in CODE_EXTS or ext in CONFIG_EXTS:
        return _peek_text_file(path)
    if ext in CSV_EXTS:
        return _peek_csv(path)
    if ext in PDF_EXTS:
        return _peek_pdf(path)
    if ext in DOCX_EXTS:
        return _peek_docx(path)
    if ext in XLSX_EXTS:
        return _peek_xlsx(path)
    if ext in PPTX_EXTS:
        return _peek_pptx(path)
    if ext in HTML_EXTS:
        return _peek_html(path)
    if ext in IPYNB_EXTS:
        return _peek_ipynb(path)
    if ext in IMAGE_EXTS:
        return _peek_image(path)

    # Unknown extension — return a minimal peek with an error so decide() can skip.
    size = 0
    try:
        size = path.stat().st_size
    except OSError:
        pass
    return PeekResult(
        path=path,
        ext=ext,
        size_bytes=size,
        page_count=0,
        text_density=0,
        extractable=False,
        image_ratio=0.0,
        row_count=None,
        image_dims=None,
        peek_text="",
        peek_error=f"unsupported extension: {ext or '(none)'}",
    )


# ---------------------------------------------------------------------------
# Scoring
# ---------------------------------------------------------------------------

def compute_visual_score(p: PeekResult) -> int:
    """How much does this file need the visual lane? 0–10, higher = more visual."""
    score = 0

    # PDF / DOCX text density signal
    if p.ext in PDF_EXTS:
        if p.text_density < 50:
            score += 5
        elif p.text_density < 500:
            score += 2
        if not p.extractable:
            score += 3
    if p.ext in DOCX_EXTS and p.image_ratio > 0.3:
        score += 2

    # Pure-image files are visual by nature
    if p.ext in IMAGE_EXTS and p.image_dims is not None:
        w, h = p.image_dims
        if w >= IMAGE_THUMBNAIL_MIN_DIM and h >= IMAGE_THUMBNAIL_MIN_DIM:
            score += 4
        aspect = w / h if h else 1.0
        if aspect < 0.65 or aspect > 1.6:
            score += 1

    return min(score, 10)


# Sensitivity detectors — content-only. Filenames are intentionally unused.
_RE_CURRENCY = re.compile(r"[\$€£¥]\s?\d+[.,]\d{2}")
_RE_TOTALS = re.compile(
    r"\b(total|subtotal|balance|amount\s*due|grand\s*total)\b", re.IGNORECASE
)
_RE_MASKED_ACCT = re.compile(r"(?:\*{3,}|[xX]{3,})\d{2,4}")
_RE_DATE_AMOUNT = re.compile(
    r"(\d{4}-\d{2}-\d{2}|\d{1,2}/\d{1,2})[^\n]{0,50}?\d+[.,]\d{2}"
)
_RE_LEGAL = re.compile(
    r"\b(hereby|whereas|witnesseth|effective\s+date|in\s+witness)\b", re.IGNORECASE
)
_RE_IDS = re.compile(
    r"\b(passport\s*no|license\s*(no|number)|SSN|EIN|tax\s*id)\b", re.IGNORECASE
)


def compute_sensitivity_score(p: PeekResult) -> int:
    """Return 0–10 based on content patterns in peek_text. 0 if no text to scan."""
    if not p.peek_text:
        return 0

    text = p.peek_text
    score = 0
    if len(_RE_CURRENCY.findall(text)) >= 3:
        score += 2
    if len(_RE_TOTALS.findall(text)) >= 2:
        score += 2
    if _RE_MASKED_ACCT.search(text):
        score += 3
    if len(_RE_DATE_AMOUNT.findall(text)) >= 3:
        score += 3
    if _RE_LEGAL.search(text):
        score += 2
    if _RE_IDS.search(text):
        score += 3
    return min(score, 10)


def estimate_t4_cost(p: PeekResult, *, gpu_available: bool) -> tuple[float, float]:
    """Return (storage_mb, wall_clock_seconds) for running T4 on this file.

    For images, treated as a single page.
    """
    pages = max(p.page_count, 1)
    storage_mb = pages * T4_STORAGE_MB_PER_PAGE
    rate = T4_GPU_SECONDS_PER_PAGE if gpu_available else T4_CPU_SECONDS_PER_PAGE
    seconds = pages * rate
    return storage_mb, seconds


# ---------------------------------------------------------------------------
# .nasconfig.yaml — folder-level overrides
# ---------------------------------------------------------------------------

_NASCONFIG_FILENAME = ".nasconfig.yaml"


def load_nasconfig(path: Path, *, stop_at: Path | None = None) -> dict:
    """Walk parent folders looking for `.nasconfig.yaml`. First match wins (nearest folder).

    Returns a dict with any of: {"accuracy": "critical"|"normal"|"casual",
    "t4_budget_gb_override": float, "colpali": "always"|"never"}.
    Missing or invalid file → empty dict (no override).
    """
    import yaml  # deferred — small but keep import-time light

    root = path.resolve().parent if path.is_file() else path.resolve()
    stop = stop_at.resolve() if stop_at else None

    current = root
    while True:
        candidate = current / _NASCONFIG_FILENAME
        if candidate.is_file():
            try:
                with candidate.open(encoding="utf-8") as f:
                    data = yaml.safe_load(f) or {}
                if isinstance(data, dict):
                    return data
                return {}
            except (OSError, yaml.YAMLError):
                return {}
        if stop is not None and current == stop:
            break
        parent = current.parent
        if parent == current:
            break
        current = parent
    return {}


# ---------------------------------------------------------------------------
# Decision — the central policy function
# ---------------------------------------------------------------------------

def _resolve_criticality(
    sensitivity_score: int,
    nasconfig: dict,
) -> tuple[Criticality, CriticalitySource]:
    """Criticality can be upgraded by user config or auto-detection, never downgraded."""
    user_level = str(nasconfig.get("accuracy", "")).lower()

    # Auto signal: sensitivity_score ≥ 4 → critical
    auto_critical = sensitivity_score >= 4

    if user_level == "critical":
        return "critical", "user"
    if auto_critical:
        return "critical", "auto"
    if user_level == "casual":
        # "casual" cannot suppress auto-critical (handled above); otherwise keep normal.
        return "casual", "user"
    if user_level == "normal":
        return "normal", "user"
    return "normal", "default"


def _skip_decision(reason: str, peek: PeekResult) -> RouteDecision:
    return RouteDecision(
        routes=[],
        visual_score=compute_visual_score(peek),
        sensitivity_score=compute_sensitivity_score(peek),
        t4_cost_mb=0.0,
        t4_cost_s=0.0,
        criticality="normal",
        criticality_source="default",
        skip_reason=reason,
        notes=[],
    )


def decide(
    p: PeekResult,
    *,
    nasconfig: dict | None = None,
    gpu_available: bool = False,
    t4_budget_used_mb: float = 0.0,
    t4_budget_cap_mb: float | None = None,
) -> RouteDecision:
    """Route one file. Pure function of its peek, folder config, and budget state."""

    nasconfig = nasconfig or {}
    notes: list[str] = []

    if p.peek_error and not p.peek_text and p.size_bytes == 0:
        return _skip_decision(f"peek failed: {p.peek_error}", p)

    vs = compute_visual_score(p)
    ss = compute_sensitivity_score(p)
    t4_mb, t4_s = estimate_t4_cost(p, gpu_available=gpu_available)
    criticality, crit_source = _resolve_criticality(ss, nasconfig)

    # Resolve corpus T4 budget (user can override)
    budget_cap = t4_budget_cap_mb
    if budget_cap is None:
        override = nasconfig.get("t4_budget_gb_override")
        if isinstance(override, (int, float)) and override > 0:
            budget_cap = float(override) * 1024
        else:
            budget_cap = float(DEFAULT_T4_BUDGET_MB)

    # Global opt-out of ColPali in this folder
    colpali_pref = str(nasconfig.get("colpali", "")).lower()
    colpali_disabled = colpali_pref == "never"
    colpali_forced = colpali_pref == "always"

    ext = p.ext

    # --- Skip cases ---------------------------------------------------------

    if p.peek_error and p.ext not in PDF_EXTS:
        # For non-PDFs, peek errors are fatal (PDF peek can return an error string
        # but still have a valid page_count of 0 which we treat as degenerate).
        return _skip_decision(f"peek failed: {p.peek_error}", p)

    if ext in IMAGE_EXTS:
        # Thumbnail filter — small file AND small dims → skip
        small_bytes = p.size_bytes < IMAGE_THUMBNAIL_SIZE_BYTES
        small_dims = (
            p.image_dims is not None
            and p.image_dims[0] < IMAGE_THUMBNAIL_MIN_DIM
            and p.image_dims[1] < IMAGE_THUMBNAIL_MIN_DIM
        )
        if small_bytes and small_dims:
            notes.append("thumbnail skipped (bytes + dims both small)")
            return _skip_decision("thumbnail", p)

    # --- Pure-text / code / config -----------------------------------------

    if ext in TEXT_EXTS:
        tier: Tier = "T0" if p.size_bytes >= TEXT_SIZE_T0_THRESHOLD else "T1"
        return RouteDecision(
            routes=[tier], visual_score=vs, sensitivity_score=ss,
            t4_cost_mb=0.0, t4_cost_s=0.0,
            criticality=criticality, criticality_source=crit_source,
            skip_reason=None,
            notes=notes + [f"text file, size={p.size_bytes}"],
        )

    if ext in CODE_EXTS:
        tier = "T0" if p.size_bytes >= CODE_SIZE_T0_THRESHOLD else "T1"
        return RouteDecision(
            routes=[tier], visual_score=vs, sensitivity_score=ss,
            t4_cost_mb=0.0, t4_cost_s=0.0,
            criticality=criticality, criticality_source=crit_source,
            skip_reason=None,
            notes=notes + [f"code file, size={p.size_bytes}"],
        )

    if ext in CONFIG_EXTS:
        tier = "T0" if p.size_bytes >= CONFIG_SIZE_T0_THRESHOLD else "T1"
        return RouteDecision(
            routes=[tier], visual_score=vs, sensitivity_score=ss,
            t4_cost_mb=0.0, t4_cost_s=0.0,
            criticality=criticality, criticality_source=crit_source,
            skip_reason=None,
            notes=notes + [f"config file, size={p.size_bytes}"],
        )

    # --- CSV ---------------------------------------------------------------

    if ext in CSV_EXTS:
        rc = p.row_count or 0
        if rc <= CSV_ROWS_T1_MAX:
            tier = "T1"
            notes.append(f"csv {rc} rows → T1 (row-level embed)")
        elif rc <= CSV_ROWS_T2_MAX:
            tier = "T2"
            notes.append(f"csv {rc} rows → T2 (sample summary + rows)")
        else:
            tier = "T0"
            notes.append(f"csv {rc} rows → T0 (sample summary only)")
        return RouteDecision(
            routes=[tier], visual_score=vs, sensitivity_score=ss,
            t4_cost_mb=0.0, t4_cost_s=0.0,
            criticality=criticality, criticality_source=crit_source,
            skip_reason=None, notes=notes,
        )

    # --- PDF ---------------------------------------------------------------

    if ext in PDF_EXTS:
        if p.page_count == 0:
            return _skip_decision(f"pdf empty or unreadable: {p.peek_error}", p)

        # Short PDFs default to T3 (discriminator-heavy), regardless of visual_score.
        if p.page_count <= PDF_SHORT_PAGE_THRESHOLD and not colpali_forced:
            notes.append(f"short PDF ({p.page_count} pp) → T3")
            return RouteDecision(
                routes=["T3"], visual_score=vs, sensitivity_score=ss,
                t4_cost_mb=t4_mb, t4_cost_s=t4_s,
                criticality=criticality, criticality_source=crit_source,
                skip_reason=None, notes=notes,
            )

        # Text-native PDF path
        if vs < 7 and p.extractable and p.text_density >= 100:
            routes: list[Tier] = ["T2"]
            if criticality == "critical":
                routes = ["T3", "T2"]
                notes.append("text-native + critical → T3 + T2")
            else:
                notes.append("text-native, non-critical → T2")
            return RouteDecision(
                routes=routes, visual_score=vs, sensitivity_score=ss,
                t4_cost_mb=t4_mb, t4_cost_s=t4_s,
                criticality=criticality, criticality_source=crit_source,
                skip_reason=None, notes=notes,
            )

        # Visual PDF path — scanned, figure-heavy, or sparse text
        if colpali_disabled:
            notes.append("colpali disabled by config → T3")
            return RouteDecision(
                routes=["T3"], visual_score=vs, sensitivity_score=ss,
                t4_cost_mb=t4_mb, t4_cost_s=t4_s,
                criticality=criticality, criticality_source=crit_source,
                skip_reason=None, notes=notes,
            )

        # Check T4 cost gates
        t4_fits_per_file = (
            t4_mb <= T4_MAX_STORAGE_MB_PER_FILE
            and t4_s <= (
                T4_MAX_SECONDS_PER_FILE_GPU if gpu_available
                else T4_MAX_SECONDS_PER_FILE_CPU
            )
        )
        t4_fits_budget = (t4_budget_used_mb + t4_mb) <= budget_cap

        if t4_fits_per_file and t4_fits_budget:
            routes = ["T4"] if criticality != "critical" else ["T3", "T4"]
            notes.append(f"visual PDF → {'+'.join(routes)}")
            return RouteDecision(
                routes=routes, visual_score=vs, sensitivity_score=ss,
                t4_cost_mb=t4_mb, t4_cost_s=t4_s,
                criticality=criticality, criticality_source=crit_source,
                skip_reason=None, notes=notes,
            )

        # T4 gated off → fall back to T3 (LLM-summary handles visual content)
        reason = "over_per_file_cap" if not t4_fits_per_file else "budget_exhausted"
        notes.append(f"wanted T4, fell back to T3 ({reason})")
        return RouteDecision(
            routes=["T3"], visual_score=vs, sensitivity_score=ss,
            t4_cost_mb=t4_mb, t4_cost_s=t4_s,
            criticality=criticality, criticality_source=crit_source,
            skip_reason=None, notes=notes,
        )

    # --- DOCX --------------------------------------------------------------

    if ext in DOCX_EXTS:
        if p.peek_error:
            return _skip_decision(f"docx unreadable: {p.peek_error}", p)

        # Figure-heavy → visual path (same gates as PDF visual)
        if p.image_ratio > 0.3 and not colpali_disabled:
            t4_fits_per_file = (
                t4_mb <= T4_MAX_STORAGE_MB_PER_FILE
                and t4_s <= (
                    T4_MAX_SECONDS_PER_FILE_GPU if gpu_available
                    else T4_MAX_SECONDS_PER_FILE_CPU
                )
            )
            t4_fits_budget = (t4_budget_used_mb + t4_mb) <= budget_cap
            if t4_fits_per_file and t4_fits_budget:
                routes = ["T4"] if criticality != "critical" else ["T3", "T4"]
                notes.append(f"figure-heavy docx → {'+'.join(routes)}")
                return RouteDecision(
                    routes=routes, visual_score=vs, sensitivity_score=ss,
                    t4_cost_mb=t4_mb, t4_cost_s=t4_s,
                    criticality=criticality, criticality_source=crit_source,
                    skip_reason=None, notes=notes,
                )

        # Text-heavy docx
        routes = ["T2"] if criticality != "critical" else ["T3", "T2"]
        notes.append(f"docx text-heavy → {'+'.join(routes)}")
        return RouteDecision(
            routes=routes, visual_score=vs, sensitivity_score=ss,
            t4_cost_mb=t4_mb, t4_cost_s=t4_s,
            criticality=criticality, criticality_source=crit_source,
            skip_reason=None, notes=notes,
        )

    # --- XLSX --------------------------------------------------------------

    if ext in XLSX_EXTS:
        if p.peek_error:
            return _skip_decision(f"xlsx unreadable: {p.peek_error}", p)
        # Huge: peek_text is the only signal we have (no row_count for XLSX yet)
        if p.size_bytes > 10 * 1024 * 1024:
            notes.append("xlsx >10 MB → T0 (sample summary)")
            return RouteDecision(
                routes=["T0"], visual_score=vs, sensitivity_score=ss,
                t4_cost_mb=0.0, t4_cost_s=0.0,
                criticality=criticality, criticality_source=crit_source,
                skip_reason=None, notes=notes,
            )
        routes = ["T2"] if criticality != "critical" else ["T3", "T2"]
        notes.append(f"xlsx → {'+'.join(routes)}")
        return RouteDecision(
            routes=routes, visual_score=vs, sensitivity_score=ss,
            t4_cost_mb=0.0, t4_cost_s=0.0,
            criticality=criticality, criticality_source=crit_source,
            skip_reason=None, notes=notes,
        )

    # --- PPTX --------------------------------------------------------------

    if ext in PPTX_EXTS:
        if p.peek_error or p.page_count == 0:
            return _skip_decision(f"pptx unreadable: {p.peek_error or 'empty'}", p)

        # Image-heavy decks need T4 (ColPali) alongside T2 for visual retrieval;
        # text-heavy decks are fine with T2 alone (or T3+T2 if critical).
        image_heavy = p.image_ratio >= 0.5

        if image_heavy and not colpali_disabled:
            t4_fits_per_file = (
                t4_mb <= T4_MAX_STORAGE_MB_PER_FILE
                and t4_s <= (
                    T4_MAX_SECONDS_PER_FILE_GPU if gpu_available
                    else T4_MAX_SECONDS_PER_FILE_CPU
                )
            )
            t4_fits_budget = (t4_budget_used_mb + t4_mb) <= budget_cap
            if t4_fits_per_file and t4_fits_budget:
                routes = (["T3", "T2", "T4"]
                          if criticality == "critical" else ["T2", "T4"])
                notes.append(f"pptx image-heavy → {'+'.join(routes)}")
                return RouteDecision(
                    routes=routes, visual_score=vs, sensitivity_score=ss,
                    t4_cost_mb=t4_mb, t4_cost_s=t4_s,
                    criticality=criticality, criticality_source=crit_source,
                    skip_reason=None, notes=notes,
                )
            notes.append(
                f"pptx image-heavy, T4 gated off "
                f"({'over_per_file_cap' if not t4_fits_per_file else 'budget_exhausted'})"
            )

        routes = ["T3", "T2"] if criticality == "critical" else ["T2"]
        notes.append(f"pptx text-heavy → {'+'.join(routes)}")
        return RouteDecision(
            routes=routes, visual_score=vs, sensitivity_score=ss,
            t4_cost_mb=t4_mb, t4_cost_s=t4_s,
            criticality=criticality, criticality_source=crit_source,
            skip_reason=None, notes=notes,
        )

    # --- HTML --------------------------------------------------------------

    if ext in HTML_EXTS:
        if not p.extractable or p.text_density == 0:
            return _skip_decision("html extracted empty — likely JS-only SPA", p)
        routes = ["T3", "T2"] if criticality == "critical" else ["T2"]
        notes.append(f"html → {'+'.join(routes)}")
        return RouteDecision(
            routes=routes, visual_score=vs, sensitivity_score=ss,
            t4_cost_mb=0.0, t4_cost_s=0.0,
            criticality=criticality, criticality_source=crit_source,
            skip_reason=None, notes=notes,
        )

    # --- IPYNB -------------------------------------------------------------

    if ext in IPYNB_EXTS:
        if p.peek_error:
            return _skip_decision(f"ipynb unreadable: {p.peek_error}", p)
        if p.page_count == 0:
            return _skip_decision("ipynb has no cells", p)
        # Notebooks are "text once parsed" — same routing as text-native docs.
        routes = ["T3", "T2"] if criticality == "critical" else ["T2"]
        notes.append(f"ipynb ({p.page_count} cells) → {'+'.join(routes)}")
        return RouteDecision(
            routes=routes, visual_score=vs, sensitivity_score=ss,
            t4_cost_mb=0.0, t4_cost_s=0.0,
            criticality=criticality, criticality_source=crit_source,
            skip_reason=None, notes=notes,
        )

    # --- Images ------------------------------------------------------------

    if ext in IMAGE_EXTS:
        # Thumbnail case handled above; we reach here for normal images.
        if colpali_disabled:
            notes.append("colpali disabled for images → T3")
            return RouteDecision(
                routes=["T3"], visual_score=vs, sensitivity_score=ss,
                t4_cost_mb=t4_mb, t4_cost_s=t4_s,
                criticality=criticality, criticality_source=crit_source,
                skip_reason=None, notes=notes,
            )
        t4_fits_per_file = (
            t4_mb <= T4_MAX_STORAGE_MB_PER_FILE
            and t4_s <= (
                T4_MAX_SECONDS_PER_FILE_GPU if gpu_available
                else T4_MAX_SECONDS_PER_FILE_CPU
            )
        )
        t4_fits_budget = (t4_budget_used_mb + t4_mb) <= budget_cap
        if t4_fits_per_file and t4_fits_budget:
            routes = ["T4"] if criticality != "critical" else ["T3", "T4"]
            notes.append(f"image → {'+'.join(routes)}")
            return RouteDecision(
                routes=routes, visual_score=vs, sensitivity_score=ss,
                t4_cost_mb=t4_mb, t4_cost_s=t4_s,
                criticality=criticality, criticality_source=crit_source,
                skip_reason=None, notes=notes,
            )
        notes.append("image T4 gated off → T3")
        return RouteDecision(
            routes=["T3"], visual_score=vs, sensitivity_score=ss,
            t4_cost_mb=t4_mb, t4_cost_s=t4_s,
            criticality=criticality, criticality_source=crit_source,
            skip_reason=None, notes=notes,
        )

    # --- Unknown -----------------------------------------------------------
    return _skip_decision(f"unsupported extension: {ext or '(none)'}", p)


# ---------------------------------------------------------------------------
# CLI: `python -m src.router <file>` — inspect a single file's routing decision
# ---------------------------------------------------------------------------

def _detect_gpu() -> bool:
    """Best-effort: detect CUDA or MPS without importing torch unless present."""
    try:
        import torch  # deferred
    except ImportError:
        return False
    try:
        if torch.cuda.is_available():
            return True
        if hasattr(torch.backends, "mps") and torch.backends.mps.is_available():
            return True
    except Exception:  # pylint: disable=broad-except
        pass
    return False


def _cli_explain_file(path: Path, gpu: bool) -> int:
    """Print the full routing decision for one file."""
    peeked = peek(path)
    nasconfig = load_nasconfig(path)
    decision = decide(peeked, nasconfig=nasconfig, gpu_available=gpu)

    print(f"path:          {path}")
    print(f"ext:           {peeked.ext}")
    print(f"size_bytes:    {peeked.size_bytes:,}")
    if peeked.page_count:
        print(f"pages:         {peeked.page_count}")
    if peeked.row_count is not None:
        print(f"rows:          {peeked.row_count:,}")
    if peeked.image_dims:
        print(f"image_dims:    {peeked.image_dims[0]}x{peeked.image_dims[1]}")
    print(f"text_density:  {peeked.text_density}")
    print(f"extractable:   {peeked.extractable}")
    if peeked.peek_error:
        print(f"peek_error:    {peeked.peek_error}")
    print()
    print(f"visual_score:       {decision.visual_score}")
    print(f"sensitivity_score:  {decision.sensitivity_score}")
    print(f"criticality:        {decision.criticality}  (source: {decision.criticality_source})")
    if decision.t4_cost_mb > 0:
        print(
            f"t4_cost:            {decision.t4_cost_mb:.1f} MB, "
            f"{decision.t4_cost_s:.0f}s {'(gpu)' if gpu else '(cpu)'}"
        )
    print(f"gpu_detected:       {'yes' if gpu else 'no'}")
    print()
    if decision.skipped:
        print(f"decision:      SKIP ({decision.skip_reason})")
    else:
        print(f"decision:      run {'+'.join(decision.routes)}")
    if decision.notes:
        print("notes:")
        for n in decision.notes:
            print(f"  - {n}")
    return 0


def _cli_explain_dir(root: Path, gpu: bool, *, limit: int | None) -> int:
    """Walk `root`, print one compact line per candidate, tally by tier.

    Respects `.gitignore` / `.nasignore` + built-in defaults (same rules the
    ingest walker uses). No side effects — nothing is written, no LLM.
    """
    # Deferred imports: only hit the walker/ignore modules in directory mode.
    from src.ingest.ignore import IgnoreRules
    from src.ingest.walker import find_candidates

    files, ignored, asset_lib_skipped = find_candidates(
        root, ignore_rules=IgnoreRules.from_root(root)
    )
    if not files and ignored == 0 and asset_lib_skipped == 0:
        print(f"no indexable files under {root}")
        return 0

    tally: dict[str, int] = {"T0": 0, "T1": 0, "T2": 0, "T3": 0, "T4": 0, "SKIP": 0}
    multi_tier = 0
    shown = 0

    for f in files:
        peeked = peek(f)
        nasconfig = load_nasconfig(f)
        d = decide(peeked, nasconfig=nasconfig, gpu_available=gpu)
        if d.skipped:
            tally["SKIP"] += 1
            badge = "SKIP"
            detail = d.skip_reason or ""
        else:
            route_str = "+".join(d.routes) if d.routes else "-"
            badge = f"[{route_str}]"
            for t in d.routes:
                tally[t] = tally.get(t, 0) + 1
            if len(d.routes) > 1:
                multi_tier += 1
            detail = (
                f"visual={d.visual_score} sens={d.sensitivity_score} crit={d.criticality}"
            )
            if d.t4_cost_mb > 0:
                detail += f" t4_mb={d.t4_cost_mb:.1f}"

        try:
            rel = str(f.relative_to(root))
        except ValueError:
            rel = str(f)

        if limit is None or shown < limit:
            print(f"  {badge:<10}  {rel:<70}  ({detail})")
            shown += 1
        elif shown == limit:
            remaining = len(files) - shown
            print(f"  ... +{remaining} more (pass --limit 0 to see all)")
            shown += 1

    print()
    print(
        f"summary: {len(files)} candidates — "
        f"T0={tally['T0']} T1={tally['T1']} T2={tally['T2']} "
        f"T3={tally['T3']} T4={tally['T4']} SKIP={tally['SKIP']} "
        f"(multi-tier={multi_tier}, ignored={ignored}, "
        f"asset_lib_skipped={asset_lib_skipped}, gpu={'yes' if gpu else 'no'})"
    )
    return 0


def main() -> None:
    import argparse as _argparse
    parser = _argparse.ArgumentParser(
        prog="nas-router",
        description=(
            "Inspect what tier a file (or every file in a directory) would route to. "
            "No ingest, no LLM, no writes."
        ),
    )
    parser.add_argument("path", help="File or directory to inspect.")
    parser.add_argument(
        "--limit",
        type=int,
        default=50,
        help="For directories: max per-file lines to print (default: 50, 0 = all).",
    )
    args = parser.parse_args()

    target = Path(args.path)
    if not target.exists():
        print(f"error: no such path: {target}")
        raise SystemExit(1)

    gpu = _detect_gpu()
    limit = None if args.limit == 0 else args.limit

    if target.is_dir():
        raise SystemExit(_cli_explain_dir(target.resolve(), gpu, limit=limit))
    raise SystemExit(_cli_explain_file(target, gpu))


if __name__ == "__main__":
    main()
