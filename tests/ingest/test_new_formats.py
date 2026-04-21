"""Tests for the four new file types: .log, .pptx, .html, .ipynb.

Covers:
  * Extractors in src/content.py
  * Peek + decide in src/router.py
  * tier2._extract dispatch
  * End-to-end via tier2.run → summary markdown round-trips through the Stage-2 parser
"""

from __future__ import annotations

import json
from pathlib import Path

import pytest

from src.content import (
    extract_html_text,
    extract_ipynb_text,
    extract_pptx_text,
    SummarizeError,
)
from src.ingest import tier2
from src.router import decide, peek
from src.stage2.parser import parse_summary_file


# ---------------------------------------------------------------------------
# Fixture builders
# ---------------------------------------------------------------------------

def _make_pptx(tmp_path: Path, *, slides: list[tuple[str, list[str]]],
               notes: list[str] | None = None) -> Path:
    """Build a .pptx with given slide titles + bullets."""
    from pptx import Presentation
    prs = Presentation()
    blank = prs.slide_layouts[6]           # fully blank layout
    for i, (title, bullets) in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        # Title textbox
        tb = slide.shapes.add_textbox(left=914_400, top=914_400,
                                      width=6_858_000, height=914_400)
        tb.text_frame.text = title
        # Body bullets textbox
        if bullets:
            body = slide.shapes.add_textbox(left=914_400, top=1_828_800,
                                            width=6_858_000, height=3_200_000)
            body.text_frame.text = "\n".join(bullets)
        if notes and i < len(notes) and notes[i]:
            slide.notes_slide.notes_text_frame.text = notes[i]
    out = tmp_path / "deck.pptx"
    prs.save(str(out))
    return out


def _make_html(tmp_path: Path, *, article: str, include_chrome: bool = True) -> Path:
    """Build an HTML page with boilerplate around the real article."""
    chrome = (
        "<html><head><title>X</title><script>alert(1)</script>"
        "<style>body{color:red}</style></head><body>"
        "<nav><a>Home</a><a>About</a></nav>"
        "<header><h1>Site name</h1></header>"
    ) if include_chrome else "<html><body>"
    footer = (
        "<footer><p>Copyright 2026 Site Inc.</p></footer>"
        "</body></html>"
    ) if include_chrome else "</body></html>"
    full = f"{chrome}<article><h1>Article</h1><p>{article}</p></article>{footer}"
    p = tmp_path / "page.html"
    p.write_text(full, encoding="utf-8")
    return p


def _make_ipynb(tmp_path: Path, cells: list[dict]) -> Path:
    nb = {
        "nbformat": 4,
        "nbformat_minor": 5,
        "metadata": {"kernelspec": {"name": "python3", "display_name": "Python 3"}},
        "cells": cells,
    }
    p = tmp_path / "notebook.ipynb"
    p.write_text(json.dumps(nb), encoding="utf-8")
    return p


# ---------------------------------------------------------------------------
# .log
# ---------------------------------------------------------------------------

def test_log_small_routes_to_t1(tmp_path: Path):
    p = tmp_path / "app.log"
    p.write_text("INFO started\nERROR connection lost\n", encoding="utf-8")
    d = decide(peek(p))
    assert d.routes == ["T1"]


def test_log_huge_routes_to_t0(tmp_path: Path):
    p = tmp_path / "server.log"
    # >100 KB threshold → T0
    p.write_text("ERROR db timeout\n" * 20_000, encoding="utf-8")
    d = decide(peek(p))
    assert d.routes == ["T0"]


# ---------------------------------------------------------------------------
# .pptx
# ---------------------------------------------------------------------------

def test_pptx_text_heavy_routes_to_t2(tmp_path: Path):
    d_pptx = _make_pptx(tmp_path, slides=[
        ("Intro", ["Bullet one", "Bullet two"]),
        ("Body", ["Point A", "Point B", "Point C"]),
        ("Conclusion", ["Summary"]),
    ])
    d = decide(peek(d_pptx), gpu_available=True)
    assert d.routes == ["T2"]


def test_pptx_extract_pulls_title_bullets_notes(tmp_path: Path):
    p = _make_pptx(
        tmp_path,
        slides=[("Beethoven's 5th", ["Sonata form", "Key: C minor"])],
        notes=["Speaker note: emphasize the motif"],
    )
    text = extract_pptx_text(p)
    assert "Beethoven" in text
    assert "Sonata form" in text
    assert "Speaker note" in text
    assert "## Slide 1" in text


def test_pptx_round_trips_through_tier2(tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
    sdir = tmp_path / "summaries"
    sdir.mkdir()
    monkeypatch.setattr("src.ingest.common.SUMMARIES_DIR", sdir)
    monkeypatch.setattr(
        "src.ingest.tier2.summary_output_path",
        lambda path, tier: sdir / f"{path.stem}_{tier}.md",
    )

    p = _make_pptx(tmp_path, slides=[("Intro", ["First bullet"])])
    tier2.run(p, "deck.pptx")
    md = sdir / "deck_t2.md"
    parsed = parse_summary_file(md)
    assert "Intro" in parsed.summary
    assert "First bullet" in parsed.summary
    assert parsed.content_type == "pptx"


# ---------------------------------------------------------------------------
# .html
# ---------------------------------------------------------------------------

def test_html_strips_boilerplate(tmp_path: Path):
    p = _make_html(tmp_path, article="The real content that should survive extraction.")
    text = extract_html_text(p)
    # Real article text kept
    assert "real content" in text
    # Scripts and styles MUST be gone (trafilatura's main job on this axis)
    assert "alert(1)" not in text
    assert "body{color:red}" not in text
    # Copyright footer boilerplate should be stripped
    assert "Copyright 2026" not in text


def test_html_routes_to_t2(tmp_path: Path):
    p = _make_html(tmp_path, article="Article body with several meaningful words here.")
    d = decide(peek(p))
    assert d.routes == ["T2"]


def test_html_empty_extraction_is_skipped(tmp_path: Path):
    # SPA-style HTML: no article content, just a JS placeholder.
    p = tmp_path / "spa.html"
    p.write_text(
        "<html><body><div id='root'></div><script>app.init()</script></body></html>",
        encoding="utf-8",
    )
    d = decide(peek(p))
    assert d.skipped
    assert "JS-only" in (d.skip_reason or "") or "empty" in (d.skip_reason or "")


# ---------------------------------------------------------------------------
# .ipynb
# ---------------------------------------------------------------------------

def test_ipynb_extracts_code_and_markdown(tmp_path: Path):
    p = _make_ipynb(tmp_path, cells=[
        {"cell_type": "markdown", "source": "# My analysis\n\nIntro text."},
        {"cell_type": "code", "source": "import pandas as pd\ndf = pd.read_csv('x.csv')"},
        {"cell_type": "code", "source": "df.head()", "outputs": [
            {"output_type": "stream", "text": "noise that should be skipped"}
        ]},
    ])
    text = extract_ipynb_text(p)
    assert "My analysis" in text
    assert "pandas" in text
    assert "df.head()" in text
    # Outputs not included
    assert "noise that should be skipped" not in text


def test_ipynb_routes_to_t2(tmp_path: Path):
    p = _make_ipynb(tmp_path, cells=[
        {"cell_type": "markdown", "source": "# Title\n\nBody text."},
        {"cell_type": "code", "source": "print('hi')"},
    ])
    d = decide(peek(p))
    assert d.routes == ["T2"]


def test_ipynb_malformed_json_skipped(tmp_path: Path):
    p = tmp_path / "bad.ipynb"
    p.write_text("{not json", encoding="utf-8")
    d = decide(peek(p))
    assert d.skipped


def test_ipynb_source_as_list_of_strings(tmp_path: Path):
    """Jupyter sometimes stores `source` as a list of strings; extractor must flatten."""
    p = _make_ipynb(tmp_path, cells=[
        {"cell_type": "code", "source": ["import numpy as np\n", "x = np.array([1, 2, 3])"]},
    ])
    text = extract_ipynb_text(p)
    assert "numpy" in text
    assert "np.array" in text


# ---------------------------------------------------------------------------
# Error behavior
# ---------------------------------------------------------------------------

def test_extract_pptx_raises_on_non_pptx(tmp_path: Path):
    p = tmp_path / "fake.pptx"
    p.write_bytes(b"not a real pptx")
    with pytest.raises(SummarizeError):
        extract_pptx_text(p)


def test_extract_html_raises_on_empty_file(tmp_path: Path):
    p = tmp_path / "empty.html"
    p.write_text("", encoding="utf-8")
    with pytest.raises(SummarizeError):
        extract_html_text(p)


def test_extract_ipynb_empty_when_no_cells(tmp_path: Path):
    """Matches the pattern of extract_docx/pdf: returns empty string when no content.
    The caller (build_content_blocks / tier2) decides whether empty is fatal."""
    p = tmp_path / "no_cells.ipynb"
    p.write_text(json.dumps({"nbformat": 4, "metadata": {}}), encoding="utf-8")
    assert extract_ipynb_text(p) == ""
